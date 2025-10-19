# veco.py
# -----------------------------------------------------------------------------
# VeCo – Stand-alone Vektorisierer & RAG-Retrieval
#
# - JSON als Fallback-Speicher (keine Installation nötig)
# - Optional: externe Storages (SQLite/Mongo) via separater "storages.py"
# - Fixes/Features:
#     * Korrekte FAISS-IDs (keine 0er-IDs mehr)
#     * Chunking mit Overlap (≈ 500–700 Tokens, char-basiert, satzsensitiv)
#     * Summaries werden zusätzlich gespeichert (nicht als Embedding-Basis)
#     * RAG-Query mit Ollama: query(database, frage, llm_model, ...)
#     * Saubere, relative Pfadbehandlung (keine absoluten User-Pfade)
#     * OPTIONAL: Speaker Diarization (via externem Modul veco_diarization.py)
#     * OPTIONAL: CNN Bild-Erkennung (torchvision) + Caption via externem Modul veco_pic_describe
#     * AUTO-Heuristik: maximale sinnvolle Pipeline je Dateityp, falls nichts explizit gesetzt
# -----------------------------------------------------------------------------

from __future__ import annotations
import os
import sys
import time
import json
import logging
import threading
import importlib
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional

import numpy as np
import torch
import whisper
from sentence_transformers import SentenceTransformer
from faiss import IndexFlatL2, IndexIDMap
import pdfplumber
import pytesseract
from PIL import Image
from docx import Document
from pptx import Presentation
from moviepy import VideoFileClip

# Optional: Ollama (nur wenn vorhanden/gewünscht)
try:
    import ollama
except Exception:
    ollama = None

# Optional: Vision – torchvision Klassifikation (ResNet50)
try:
    import torchvision
    from torchvision import transforms
    _VISION_OK = True
except Exception:
    _VISION_OK = False
    torchvision = None
    transforms = None

# ---------------------------------- Logging ----------------------------------
logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger("veco")


# ---------------------------------- Spinner ----------------------------------
class Spinner:
    """Kleiner CLI-Spinner für längere Operationen (nicht kritisch, aber nett)."""
    spinner_cycle = ["|", "/", "-", "\\"]

    def __init__(self, message: str = "Processing"):
        self.stop_running = False
        self.thread = threading.Thread(target=self._run, daemon=True)
        self.message = message

    def _run(self):
        i = 0
        while not self.stop_running:
            sys.stdout.write(f"\r{self.message} {self.spinner_cycle[i % len(self.spinner_cycle)]}")
            sys.stdout.flush()
            time.sleep(0.08)
            i += 1
        sys.stdout.write("\r" + " " * (len(self.message) + 2) + "\r")

    def start(self):
        try:
            self.thread.start()
        except RuntimeError:
            pass

    def stop(self):
        self.stop_running = True
        try:
            self.thread.join(timeout=0.2)
        except RuntimeError:
            pass


# ---------------------------- Hilfen / Utilities -----------------------------
def _try_import_storages():
    """Lazy-Import eines optionalen Moduls 'storages.py' (SQLite/Mongo-Backends)."""
    try:
        return importlib.import_module("storages")
    except Exception:
        return None


def _try_import_diarization():
    """Lazy-Import eines optionalen Moduls 'veco_diarization.py' (Speaker Diarization)."""
    try:
        return importlib.import_module("veco_diarization")
    except Exception:
        return None


def _try_import_pic_describe():
    """Lazy-Import eines optionalen Moduls 'veco_pic_describe' (externe Bildbeschreibung)."""
    try:
        return importlib.import_module("veco_pic_describe")
    except Exception:
        return None


def _relpath(p: str) -> str:
    """Gibt Pfad p möglichst relativ zum aktuellen Arbeitsverzeichnis zurück."""
    try:
        if os.path.isabs(p):
            return os.path.relpath(p, start=os.getcwd())
        return p
    except Exception:
        return p  # Bei Problemen unverändert zurückgeben


def chunk_text(text: str, chunk_chars: int = 1800, overlap_chars: int = 200) -> List[str]:
    """
    Char-basiertes Chunking mit Overlap (ohne zusätzliche Abhängigkeiten).
    - Versucht, bevorzugt an Satzenden ('.') zu schneiden.
    - Richtwert: 1800 Zeichen ≈ 500–700 Tokens (je nach Sprache/Text).
    """
    text = (text or "").strip()
    if not text:
        return []
    chunks: List[str] = []
    n = len(text)
    i = 0
    while i < n:
        end = min(i + chunk_chars, n)
        cut = text.rfind(".", i, end)
        if cut == -1 or cut < i + int(0.6 * chunk_chars):
            cut = end
        chunk = text[i:cut].strip()
        if chunk:
            chunks.append(chunk)
        if cut >= n:
            break
        i = max(0, cut - overlap_chars)
    return chunks


# --------------------------- Hauptklasse: Vectorize ---------------------------
class Vectorize:
    """
    - Lädt/initialisiert Embedding- & ASR-Modelle
    - Extrahiert Text aus Dateien (txt/pdf/docx/pptx/image/audio/video)
    - Chunked, vektorisiert und indexiert (FAISS)
    - Speichert Daten in JSON (Fallback) oder optionalen Backends (SQLite/Mongo)
    - OPTIONAL: Speaker Diarization (externe Pipeline), Bild-Klassifikation (torchvision),
                Bildbeschreibung (externes Modul veco_pic_describe)
    - Bietet RAG-Retrieval + Ollama-basiertes Querying
    """

    def __init__(
        self,
        default_model: str = "gemma3:12b",
        preload_json_path: Optional[str] = "vector_db.json",
        storage: Optional[object] = None,
        storage_kind: Optional[str] = None,
        storage_kwargs: Optional[dict] = None,
        write_through: bool = True,
    ):
        # Basiskonfiguration
        self.default_model = default_model
        self.preload_json_path = _relpath(preload_json_path or "vector_db.json")
        self.write_through = write_through

        # Interner Speicher (für Fallback & schnelles Save/Load)
        self.outputdb: List[Dict[str, Any]] = []
        self.id_lookup: Dict[int, Dict[str, Any]] = {}

        # Optionales externes Storage (SQLite/Mongo)
        self._ext_storage = None
        if storage is not None:
            self._ext_storage = storage
        elif storage_kind is not None:
            _stor = _try_import_storages()
            if _stor is not None:
                if storage_kind.lower() == "sqlite":
                    self._ext_storage = _stor.SqliteStorage(**(storage_kwargs or {}))
                elif storage_kind.lower() == "mongo":
                    self._ext_storage = _stor.MongoStorage(**(storage_kwargs or {}))
                else:
                    raise ValueError(f"Unbekanntes storage_kind: {storage_kind}")
            else:
                logger.warning("storages.py nicht gefunden – bleibe beim JSON-Fallback.")

        # Modelle laden (Whisper + SBERT, optional Ollama-Check)
        spinner = Spinner("Initializing models")
        spinner.start()
        try:
            device = "cuda" if torch.cuda.is_available() else "cpu"
            self.whisper_model = whisper.load_model("base", device=device)
            self.embedder = SentenceTransformer("all-MiniLM-L6-v2")
            emb_dim = self.embedder.get_sentence_embedding_dimension()
            self.faiss_index = IndexIDMap(IndexFlatL2(emb_dim))
            if ollama is not None:
                self.check_ollama_models()
        finally:
            spinner.stop()

        # Vision (optional) – Lazy Init Handles
        self._vision_cls = None  # torchvision resnet50
        self._vision_tf = None   # zugehörige Transforms

        # Bootstrap: vorhandene Daten laden (Storage bevorzugt, sonst JSON)
        if self._ext_storage is not None:
            self._bootstrap_from_storage()
        else:
            self.load_database()

    # ---------------------- Infrastruktur / I/O ------------------------
    def check_ollama_models(self):
        """Nur zum Prüfen, ob Ollama erreichbar ist – kein harter Fehler."""
        try:
            _ = ollama.list()
        except Exception:
            logger.info("Ollama nicht verfügbar – LLM-Kompression/RAG-Antwort ggf. deaktiviert.")

    def detect_input_type(self, path: str) -> str:
        """Dateityp heuristisch per Endung bestimmen."""
        p = str(path).lower()
        if p.endswith(".txt"):
            return "text"
        if p.endswith(".pdf"):
            return "pdf"
        if p.endswith(".docx"):
            return "word"
        if p.endswith(".pptx"):
            return "pptx"
        if p.endswith((".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff")):
            return "image"
        if p.endswith((".wav", ".mp3", ".m4a", ".flac", ".ogg")):
            return "audio"
        if p.endswith((".mp4", ".mov", ".mkv", ".avi")):
            return "video"
        return "text"

    # -------------------------- Extraktion ----------------------------
    def extract_text(self, inputfile: str, input_type: str) -> str:
        """Extrahiert Text aus Text/PDF/Word/PPTX-Dateien."""
        if input_type == "text":
            return Path(inputfile).read_text(encoding="utf-8", errors="ignore")

        if input_type == "pdf":
            texts = []
            with pdfplumber.open(inputfile) as pdf:
                for page in pdf.pages:
                    t = page.extract_text() or ""
                    texts.append(t)
            return "\n".join(texts)

        if input_type == "word":
            doc = Document(inputfile)
            return "\n".join(p.text for p in doc.paragraphs)

        if input_type == "pptx":
            prs = Presentation(inputfile)
            slides = []
            for slide in prs.slides:
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
                slides.append("\n".join(parts))
            return "\n\n".join(slides)

        return ""

    def extract_text_from_image(self, inputfile: str) -> str:
        """OCR via Tesseract (deu+eng)."""
        img = Image.open(inputfile)
        txt = pytesseract.image_to_string(img, lang="deu+eng")
        return txt

    def extract_text_from_audio(self, inputfile: str) -> str:
        """ASR via Whisper-Modell (erzwungen Deutsch; bei Bedarf anpassbar)."""
        result = self.whisper_model.transcribe(inputfile, language="de")
        return result.get("text", "")

    def extract_text_from_video(self, inputfile: str) -> str:
        """
        Einfache Videoverarbeitung:
        - Audio extrahieren -> temporäre WAV im aktuellen Ordner
        - Whisper-Transkription
        """
        base = Path(inputfile).stem
        tmp_wav = f"{base}.veco_tmp.wav"  # relativ, kein User-Home-Pfad
        clip = VideoFileClip(inputfile)
        clip.audio.write_audiofile(tmp_wav, verbose=False, logger=None)
        text = self.extract_text_from_audio(tmp_wav)
        try:
            os.remove(tmp_wav)
        except Exception:
            pass
        return text

    # ----------------------- OPTIONAL: Diarization ---------------------
    def _run_diarization(self, inputfile: str, diarization_kwargs: Optional[dict] = None) -> Optional[str]:
        """
        Führt (wenn verfügbar) deine externe Diarization-Pipeline aus:
        - erwartet ein Modul 'veco_diarization.py' mit 'run_file' und 'build_config'
        - gibt einen Speaker-getaggten Text zurück oder None bei Fehler/fehlendem Modul
        """
        dia = _try_import_diarization()
        if dia is None:
            logger.info("Diarization-Modul (veco_diarization.py) nicht gefunden – überspringe.")
            return None

        with tempfile.TemporaryDirectory(prefix="veco_dia_") as tmpdir:
            try:
                kwargs = diarization_kwargs or {}
                kwargs.setdefault("audio_dir", tmpdir)
                ok, out_txt = dia.run_file(inputfile, **kwargs)
                if ok and out_txt and os.path.exists(out_txt):
                    return Path(out_txt).read_text(encoding="utf-8", errors="ignore")
            except Exception as e:
                logger.warning(f"Diarization fehlgeschlagen: {e}")
        return None

    # ----------------------- OPTIONAL: Vision/CNN ----------------------
    def _init_vision_classifier(self):
        if self._vision_cls is not None:
            return
        if not _VISION_OK:
            logger.info("torchvision nicht verfügbar – Bildklassifikation wird übersprungen.")
            return
        try:
            weights = torchvision.models.ResNet50_Weights.DEFAULT
            self._vision_cls = torchvision.models.resnet50(weights=weights)
            self._vision_cls.eval()
            self._vision_tf = weights.transforms()
        except Exception:
            # Fallback auf einfache Normalisierung/Resize
            self._vision_cls = torchvision.models.resnet50(pretrained=True)
            self._vision_cls.eval()
            self._vision_tf = transforms.Compose([
                transforms.Resize(256),
                transforms.CenterCrop(224),
                transforms.ToTensor(),
                transforms.Normalize(mean=[0.485, 0.456, 0.406],
                                     std =[0.229, 0.224, 0.225]),
            ])

    def _image_classify(self, inputfile: str, topk: int = 5) -> Optional[str]:
        self._init_vision_classifier()
        if self._vision_cls is None or self._vision_tf is None:
            return None
        try:
            img = Image.open(inputfile).convert("RGB")
            x = self._vision_tf(img).unsqueeze(0)
            with torch.no_grad():
                logits = self._vision_cls(x)
                probs = torch.softmax(logits, dim=1)[0]
                k = min(topk, probs.shape[0])
                topv, topi = torch.topk(probs, k=k)
                try:
                    labels = torchvision.models.ResNet50_Weights.DEFAULT.meta["categories"]
                except Exception:
                    labels = [f"class_{i}" for i in range(probs.shape[0])]
                pairs = [f"{labels[int(i)]} ({float(v):.2%})" for v, i in zip(topv, topi)]
                return "Bildklassifikation (Top): " + ", ".join(pairs)
        except Exception as e:
            logger.warning(f"Bildklassifikation fehlgeschlagen: {e}")
            return None

    def _image_caption_external(self, inputfile: str, **kwargs) -> Optional[str]:
        """
        Nutzt ein externes Projekt 'veco_pic_describe' (falls vorhanden).
        Erwartete API (mindestens eine der folgenden):
          - describe(image_path: str, **kwargs) -> str
          - run(image_path: str, **kwargs) -> str
        """
        mod = _try_import_pic_describe()
        if mod is None:
            logger.info("veco_pic_describe nicht gefunden – Bildbeschreibung wird übersprungen.")
            return None
        try:
            if hasattr(mod, "describe") and callable(mod.describe):
                text = mod.describe(inputfile, **kwargs)
            elif hasattr(mod, "run") and callable(mod.run):
                text = mod.run(inputfile, **kwargs)
            else:
                logger.warning("veco_pic_describe hat weder describe() noch run().")
                return None
            text = (text or "").strip()
            return f"Bildbeschreibung: {text}" if text else None
        except Exception as e:
            logger.warning(f"Externe Bildbeschreibung fehlgeschlagen: {e}")
            return None

    # ----------------------- Embedding & Index -------------------------
    def embed_texts(self, texts: List[str]) -> np.ndarray:
        """Batch-Embedding für eine Liste von Strings (returns float32 ndarray)."""
        if not texts:
            dim = self.embedder.get_sentence_embedding_dimension()
            return np.zeros((0, dim), dtype=np.float32)
        vecs = self.embedder.encode(texts, convert_to_numpy=True, normalize_embeddings=False)
        return np.asarray(vecs, dtype=np.float32)

    def _next_id(self) -> int:
        """Vergibt neue, globale Record-IDs (stabil über JSON/Storage)."""
        if self.outputdb:
            return max(int(r["id"]) for r in self.outputdb if "id" in r) + 1
        try:
            if self._ext_storage is not None:
                return self._ext_storage.get_max_id() + 1
        except Exception:
            pass
        return 0

    def _add_records(self, vectors: np.ndarray, chunks: List[str], source: str, doc_id: int):
        """
        Fügt Chunks + Embeddings in Index & DB ein – mit korrekten FAISS-IDs.
        - vectors: (N, dim)
        - chunks:  Liste der Chunk-Texte (N)
        - source:  (relativer) Pfad zur Quelle
        - doc_id:  logische Gruppierung (erste ID dieses Dokuments)
        """
        assert vectors.shape[0] == len(chunks)
        start_id = self._next_id()
        ids = np.arange(start_id, start_id + len(chunks), dtype=np.int64)

        # FAISS: Embeddings mit expliziten IDs hinzufügen
        self.faiss_index.add_with_ids(vectors, ids)

        # Records bauen & persistieren (in-memory + optional extern)
        src = _relpath(source)
        for local_idx, (rid, chunk, vec) in enumerate(zip(ids.tolist(), chunks, vectors)):
            rec: Dict[str, Any] = {
                "id": int(rid),
                "doc_id": int(doc_id),
                "chunk_idx": local_idx,
                "text": chunk,              # Embedding-Basis = Original-Chunk
                "source": src,              # relative Quelle
                "vector": vec.tolist(),     # nützlich für Re-Index/Laden
            }
            self.outputdb.append(rec)
            self.id_lookup[int(rid)] = rec
            if self._ext_storage is not None and self.write_through:
                self._ext_storage.upsert(rec)

    # --------------------------- Persistenz ----------------------------
    def save_database(self, json_path: Optional[str] = None):
        """
        JSON-Fallback: schreibt self.outputdb in eine JSON-Datei.
        Wenn externes Storage aktiv & write_through=False: zusätzlich dorthin upserten.
        """
        path = json_path or self.preload_json_path
        payload = {"outputdb": self.outputdb}
        Path(path).write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

        if self._ext_storage is not None and not self.write_through:
            for rec in self.outputdb:
                self._ext_storage.upsert(rec)

        logger.info(f"JSON gespeichert: {_relpath(str(path))}")

    def load_database(self, json_path: Optional[str] = None):
        """
        JSON-Fallback laden und FAISS + id_lookup füllen.
        Existiert die Datei nicht, wird leer gestartet (keine Exception).
        """
        self.outputdb.clear()
        self.id_lookup.clear()

        path = json_path or self.preload_json_path
        if not Path(path).exists():
            logger.info(f"Kein JSON vorhanden ({_relpath(str(path))}) – starte leer.")
            return

        data = json.loads(Path(path).read_text(encoding="utf-8"))
        cnt = 0
        for rec in data.get("outputdb", []):
            self.outputdb.append(rec)
            rid = rec.get("id", None)
            vec = rec.get("vector", None)
            if rid is not None and isinstance(vec, list):
                v = np.array([vec], dtype=np.float32)
                self.faiss_index.add_with_ids(v, np.array([int(rid)]))
                self.id_lookup[int(rid)] = rec
                cnt += 1

        logger.info(f"JSON geladen: {cnt} Vektoren in FAISS.")

    def _bootstrap_from_storage(self):
        """
        Daten aus externem Storage (SQLite/Mongo) laden und Index/Lookup füllen.
        JSON bleibt unberührt.
        """
        self.outputdb.clear()
        self.id_lookup.clear()

        cnt = 0
        for rec in self._ext_storage.load_all():
            self.outputdb.append(rec)
            rid = rec.get("id", None)
            vec = rec.get("vector", None)
            if rid is not None and isinstance(vec, list):
                v = np.array([vec], dtype=np.float32)
                self.faiss_index.add_with_ids(v, np.array([int(rid)]))
            if rid is not None:
                self.id_lookup[int(rid)] = rec
            cnt += 1

        logger.info(f"Storage geladen: {cnt} Datensätze.")

    def _switch_database(self, database: str):
        """
        Dynamischer DB-Wechsel anhand des 'database'-Strings:
          - *.json           -> JSON-Fallback
          - *.sqlite / *.db  -> SQLite (braucht storages.SqliteStorage)
          - mongodb://...    -> MongoDB (braucht storages.MongoStorage)
        Baut danach FAISS & Lookup neu auf.
        """
        db = (database or "").strip()

        # Bereits aktiv? Dann no-op.
        if getattr(self, "_active_db", None) == db:
            return

        # Index & Caches resetten
        emb_dim = self.embedder.get_sentence_embedding_dimension()
        self.faiss_index = IndexIDMap(IndexFlatL2(emb_dim))
        self.outputdb.clear()
        self.id_lookup.clear()

        # Externes Storage ggf. schließen
        if self._ext_storage is not None:
            try:
                self._ext_storage.close()
            except Exception:
                pass
        self._ext_storage = None

        # Routing nach Schema/Endung
        lower = db.lower()
        stor_mod = _try_import_storages()

        if lower.endswith(".json") or lower == "":
            # JSON-Fallback
            self.preload_json_path = db if db else self.preload_json_path
            self.load_database(self.preload_json_path)
            self._active_db = db
            logger.info(f"Switched to JSON database: {_relpath(str(self.preload_json_path))}")
            return

        if lower.endswith(".sqlite") or lower.endswith(".db"):
            if stor_mod is None:
                raise RuntimeError("SQLite verlangt 'storages.py'. Bitte bereitstellen.")
            self._ext_storage = stor_mod.SqliteStorage(db_path=db)
            self._bootstrap_from_storage()
            self._active_db = db
            logger.info(f"Switched to SQLite database: {_relpath(db)}")
            return

        if lower.startswith("mongodb://") or lower.startswith("mongodb+srv://"):
            if stor_mod is None:
                raise RuntimeError("Mongo verlangt 'storages.py'. Bitte bereitstellen.")
            self._ext_storage = stor_mod.MongoStorage(uri=db, db_name="veco_db", collection="entries")
            self._bootstrap_from_storage()
            self._active_db = db
            logger.info(f"Switched to Mongo database (uri): {db}")
            return

        raise ValueError(f"Unbekanntes database-Format: {database}")

    def close(self):
        """Aufräumen (externes Storage sauber schließen)."""
        if self._ext_storage is not None:
            try:
                self._ext_storage.close()
            except Exception as e:
                logger.warning(f"Storage close error: {e}")

    # ------------------------ LLM / Summarization ----------------------
    def build_compression_prompt(self, text: str) -> str:
        return (
            "Fasse den folgenden Text prägnant als Executive Summary (5–8 Bullet Points) zusammen.\n\n"
            f"TEXT:\n{text}\n"
        )

    def ask_llm(self, prompt: str, model: Optional[str] = None) -> str:
        if ollama is None:
            raise RuntimeError("Ollama nicht verfügbar.")
        m = model or self.default_model
        resp = ollama.generate(model=m, prompt=prompt)
        return (resp.get("response") or "").strip()

    # ------------------------ Ingest / Vektorisierung ------------------
    def vectorize(
        self,
        inputfile: str,
        use_compression: bool = False,
        model: Optional[str] = None,

        # AUTO-Heuristik:
        use_diarization: Optional[bool] = None,   # None = AUTO, True/False = erzwingen
        diarization_kwargs: Optional[dict] = None,

        vision_mode: Optional[str] = None,        # None = AUTO (für Bilder), "classify" | "caption" | "both" | ""
        topk: int = 5,                            # für Bildklassifikation
        pic_kwargs: Optional[dict] = None,        # optionale Parameter für veco_pic_describe
    ):
        """
        Voller Ingest-Pipeline-Step:
          1) Text extrahieren (Dateityp-abhängig). Falls use_diarization=True (nur Audio/Video):
             -> Speaker-getaggte Transkription wird bevorzugt (wenn Modul vorhanden).
          2) Für Bilder optional: CNN-Klassifikation und/oder externe Bildbeschreibung.
          3) Chunking mit Overlap
          4) Embedding (nur Original-Chunks)
          5) Index + Persistenz
          6) Optional: Summarization (als separates Meta, nicht als Embedding-Basis)
        """
        spinner = Spinner("Vectorizing input")
        spinner.start()
        try:
            input_type = self.detect_input_type(inputfile)
            logger.info(f"Detected input type: {input_type}")

            raw_text = ""

            # --- 1) Volltext (inkl. AUTO-Diarization für Audio/Video) ---
            if use_diarization is None:
                use_diarization = (input_type in {"audio", "video"}) and (_try_import_diarization() is not None)

            if use_diarization and input_type in {"audio", "video"}:
                dia_text = self._run_diarization(inputfile, diarization_kwargs=diarization_kwargs)
                if dia_text:
                    raw_text = dia_text
                else:
                    raw_text = (self.extract_text_from_audio(inputfile)
                                if input_type == "audio" else
                                self.extract_text_from_video(inputfile))
            else:
                if input_type in {"text", "pdf", "word", "pptx"}:
                    raw_text = self.extract_text(inputfile, input_type)
                elif input_type == "image":
                    raw_text = self.extract_text_from_image(inputfile)  # OCR-Basis
                elif input_type == "audio":
                    raw_text = self.extract_text_from_audio(inputfile)
                elif input_type == "video":
                    raw_text = self.extract_text_from_video(inputfile)
                else:
                    raw_text = ""

            raw_text = (raw_text or "").strip()

            # --- 2) Vision-Extras (AUTO) – nur bei Bildern sinnvoll ---
            vision_extra = ""
            if input_type == "image":
                if vision_mode is None:
                    can_cls = _VISION_OK
                    can_cap = (_try_import_pic_describe() is not None)
                    if can_cls and can_cap:
                        vision_mode = "both"
                    elif can_cls:
                        vision_mode = "classify"
                    elif can_cap:
                        vision_mode = "caption"
                    else:
                        vision_mode = ""

                if vision_mode in ("classify", "both"):
                    cls = self._image_classify(inputfile, topk=topk)
                    if cls:
                        vision_extra += cls
                if vision_mode in ("caption", "both"):
                    cap = self._image_caption_external(inputfile, **(pic_kwargs or {}))
                    if cap:
                        if vision_extra:
                            vision_extra += "\n"
                        vision_extra += cap

                if vision_extra:
                    raw_text = (raw_text + "\n\n" if raw_text else "") + f"[VISION]\n{vision_extra}"

            if not raw_text:
                logger.warning("Kein Text extrahiert.")
                return

            # 3) Chunking
            chunks = chunk_text(raw_text, chunk_chars=1800, overlap_chars=200)
            if not chunks:
                chunks = [raw_text]

            # 4) Embedding (nur Original-Chunks)
            vectors = self.embed_texts(chunks)

            # 5) Doc-ID (Marker für Gruppierung; unabhängig von FAISS-IDs)
            doc_id = self._next_id()

            # 6) Index + DB
            self._add_records(vectors, chunks, source=str(inputfile), doc_id=doc_id)

            # 7) Optional: Summary (NICHT als Embedding-Basis)
            if use_compression:
                try:
                    summary = self.ask_llm(self.build_compression_prompt(raw_text), model or self.default_model)
                except Exception as e:
                    logger.warning(f"Summarization failed: {e}")
                    summary = None

                if summary:
                    meta = {
                        "id": int(10_000_000_000 + doc_id),  # große ID außerhalb FAISS
                        "doc_id": int(doc_id),
                        "chunk_idx": -1,
                        "kind": "doc_summary",
                        "text": "",
                        "summary": summary,
                        "source": _relpath(str(inputfile)),
                    }
                    self.outputdb.append(meta)
                    if self._ext_storage is not None and self.write_through:
                        self._ext_storage.upsert(meta)
        finally:
            spinner.stop()

    # --------------------------- Retrieval / RAG -----------------------
    def retrieve_context(self, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """Sucht die top_k ähnlichsten Chunks zu einer Query (Embedding + FAISS)."""
        qv = self.embed_texts([query])
        D, I = self.faiss_index.search(qv, top_k)
        hits: List[Dict[str, Any]] = []
        for rid in I[0].tolist():
            if rid == -1:
                continue
            rec = self.id_lookup.get(int(rid))
            if rec:
                hits.append(rec)
        return hits

    def query_with_context(self, question: str, top_k: int = 5, include_summary: bool = True) -> Dict[str, Any]:
        """
        Liefert Kontext-Chunks + (optional) Doc-Summaries – ohne LLM-Antwort.
        Nützlich zum Debuggen der Retrieval-Ergebnisse.
        """
        ctx = self.retrieve_context(question, top_k=top_k)
        response: Dict[str, Any] = {"question": question, "contexts": ctx}
        if include_summary and ctx:
            doc_ids = list({c.get("doc_id") for c in ctx if c.get("doc_id") is not None})
            summaries = [
                rec
                for rec in self.outputdb
                if rec.get("kind") == "doc_summary" and rec.get("doc_id") in doc_ids
            ]
            if summaries:
                response["summaries"] = summaries
        return response

    def _build_rag_prompt(
        self,
        frage: str,
        contexts: List[Dict[str, Any]],
        summaries: Optional[List[Dict[str, Any]]] = None,
    ) -> str:
        """Baut einen robusten, deutschsprachigen RAG-Prompt für Ollama."""
        ctx_text = "\n\n".join(c.get("text", "") for c in contexts if c.get("text"))
        sum_text = ""
        if summaries:
            only = [s.get("summary", "") for s in summaries if s.get("summary")]
            if only:
                sum_text = "\n\nZUSAMMENFASSUNG (Dokumentebene):\n" + "\n".join(only)

        prompt = f"""Beantworte die Frage ausschließlich auf Basis des folgenden Kontexts.
Wenn die Antwort nicht eindeutig im Kontext steht, sage: "Im gegebenen Kontext nicht enthalten."

KONTEXT:
{ctx_text}
{sum_text}

FRAGE:
{frage}

ANTWORT (präzise, deutsch):
"""
        return prompt

    def query(
        self,
        database: str,
        frage: str,
        llm_model: str,
        top_k: int = 5,
        include_summary: bool = True,
    ) -> Dict[str, Any]:
        """
        High-level RAG-Query (End-to-End):
          - database: Pfad zu JSON / SQLite (.sqlite/.db) oder Mongo-URI (mongodb://…)
          - frage:    Nutzerfrage
          - llm_model: Ollama-Modellname (z. B. "gemma3:12b")
          - top_k:    Anzahl der Kontext-Chunks
        """
        if ollama is None:
            raise RuntimeError("Ollama nicht verfügbar – bitte installieren/konfigurieren.")

        # 1) Datenbank wählen/laden
        self._switch_database(database)

        # 2) Kontext
        contexts = self.retrieve_context(frage, top_k=top_k)

        # 3) Summaries (optional)
        summaries: List[Dict[str, Any]] = []
        if include_summary and contexts:
            doc_ids = list({c.get("doc_id") for c in contexts if c.get("doc_id") is not None})
            summaries = [
                rec
                for rec in self.outputdb
                if rec.get("kind") == "doc_summary" and rec.get("doc_id") in doc_ids
            ]

        # 4) Prompt bauen und LLM fragen
        prompt = self._build_rag_prompt(frage, contexts, summaries)
        try:
            resp = ollama.generate(model=llm_model, prompt=prompt)
            answer = (resp.get("response") or "").strip()
        except Exception as e:
            raise RuntimeError(f"Ollama-Fehler: {e}")

        # 5) Ergebnis zurückgeben (mit Quellenliste)
        result = {
            "question": frage,
            "model": llm_model,
            "answer": answer,
            "contexts": contexts,
            "sources": list({c.get("source") for c in contexts if c.get("source")}),
        }
        if include_summary and summaries:
            result["summaries"] = summaries
        return result


# ------------------------------ CLI / Demo -----------------------------------
if __name__ == "__main__":
    import argparse

    ap = argparse.ArgumentParser(description="VeCo – Vectorize & RAG Retrieval")
    ap.add_argument("input", nargs="?", help="Datei (txt/pdf/docx/pptx/image/audio/video)")
    ap.add_argument("--compress", action="store_true", help="Zusammenfassung zusätzlich speichern")
    ap.add_argument("--json", default="vector_db.json", help="JSON-Fallback-Datei")
    ap.add_argument("--use-sqlite", default=None, help="Pfad zu SQLite DB (optional)")
    ap.add_argument("--use-mongo", default=None, help="Mongo URI (optional, z.B. mongodb://localhost:27017)")
    ap.add_argument("--mongo-db", default="veco_db", help="Mongo DB-Name (nur CLI-Demo)")
    ap.add_argument("--mongo-col", default="entries", help="Mongo Collection (nur CLI-Demo)")
    # Optional: explizit Vision steuern (leer string, classify, caption, both)
    ap.add_argument("--vision", default=None, help="Bildmodus: classify|caption|both|'' (None=AUTO)")
    ap.add_argument("--topk", type=int, default=5, help="Top-K Klassen für Bildklassifikation")
    # Optional: explizit Diarization steuern (None=AUTO; true/false erzwingt)
    ap.add_argument("--diarize", default=None, choices=["true", "false"], help="Diarization erzwingen (true/false). None=AUTO")

    args = ap.parse_args()

    # Optionales Storage fürs Ingest (CLI-Demo)
    storage_kind = None
    storage_kwargs = None
    if args.use_sqlite:
        storage_kind = "sqlite"
        storage_kwargs = {"db_path": args.use_sqlite}
    elif args.use_mongo:
        storage_kind = "mongo"
        storage_kwargs = {"uri": args.use_mongo, "db_name": args.mongo_db, "collection": args.mongo_col}

    veco = Vectorize(
        preload_json_path=args.json,
        storage_kind=storage_kind,
        storage_kwargs=storage_kwargs,
        write_through=True,
    )

    # CLI → Mapping auf None/Bool
    diarize_flag: Optional[bool]
    if args.diarize is None:
        diarize_flag = None  # AUTO
    else:
        diarize_flag = (args.diarize.lower() == "true")

    if args.input:
        veco.vectorize(
            args.input,
            use_compression=args.compress,
            use_diarization=diarize_flag,   # None=AUTO
            diarization_kwargs=None,         # bei Bedarf mit Parametern befüllen
            vision_mode=args.vision,         # None=AUTO
            topk=args.topk,
        )
        veco.save_database(args.json)

        # Kurzer Retrieval-Test (ohne LLM)
        res = veco.query_with_context("Worum geht es im Dokument?", top_k=5, include_summary=True)
        print(json.dumps(res, ensure_ascii=False, indent=2))
    else:
        print("Kein Input übergeben. Beispiel:")
        print("  python veco.py docs/report.pdf --compress --json vector_db.json --use-sqlite data/veco.sqlite")
        print("  python veco.py sample.wav --diarize true --json vector_db.json")
        print("  python veco.py image.jpg --vision both --json vector_db.json")

    veco.close()
